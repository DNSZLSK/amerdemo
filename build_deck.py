# -*- coding: utf-8 -*-
import copy, shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

SRC = "AMER-Proposition-Site.pptx"
shutil.copyfile(SRC, "AMER-Proposition-Site.original.pptx")

NIGHT = RGBColor(0x1B, 0x26, 0x31)
GRAY  = RGBColor(0x6B, 0x6B, 0x6B)
STONE = RGBColor(0x8B, 0x7D, 0x6B)
LINE  = RGBColor(0xDD, 0xD6, 0xCB)
SAND  = RGBColor(0xEF, 0xE7, 0xDB)

PWA_LABEL = "Une app installable (PWA)"
PWA_DESC  = ("Les visiteurs ajoutent Amer sur leur telephone comme une appli, sans passer "
             "par un store, et gardent l'agenda et les evenements a portee de main.")

PWA_LABEL = "Une app installable (PWA)"
PWA_DESC  = ("Installable sur le téléphone comme une appli, sans store. "
             "Agenda toujours à portée de main.")

TITLE5 = "Et ce site, vous pouvez le faire financer"
SUB5   = ("La création du site est une dépense éligible aux aides à la vie associative. "
          "Sur facture, elle peut être couverte jusqu'à 100%.")
LEVERS = [
    ("FDVA, fonds pour la vie associative",       "Projets et outils de communication"),
    ("Subvention Ville et Agglomération",     "Saint-Brieuc, Saint-Brieuc Armor Agglo"),
    ("Département 22, Région Bretagne",   "Patrimoine et attractivité du territoire"),
    ("Mécénat d'entreprise locale",       "Jusqu'à 60% déductible pour le donateur"),
]
BANNER5 = ("Réglé sur facture, récupérable via les aides : le site peut ne presque "
           "rien coûter à l'association.")

prs = Presentation(SRC)

# ----------------------------------------------------------------- SLIDE 2
s2 = prs.slides[1]
pre = {sp.shape_id: sp for sp in s2.shapes}
spTree = s2.shapes._spTree

def clone(src_shape, new_id, new_name, new_text=None):
    el = copy.deepcopy(src_shape._element)
    cNvPr = el.find('.//' + qn('p:cNvPr'))
    cNvPr.set('id', str(new_id)); cNvPr.set('name', new_name)
    if new_text is not None:
        ts = el.findall('.//' + qn('a:t'))
        ts[0].text = new_text
        for extra in ts[1:]:
            extra.text = ''
    spTree.append(el)

clone(pre[18], 100, "Text PWA label", PWA_LABEL)   # from "Formation admin" label
clone(pre[19], 101, "Text PWA desc",  PWA_DESC)    # from formation description
clone(pre[17], 102, "Shape PWA line")              # a divider

now = {sp.shape_id: sp for sp in s2.shapes}
labels   = [3, 6, 9, 12, 15, 100, 18]   # PWA (100) placed before Formation admin (18)
descs    = [4, 7, 10, 13, 16, 101, 19]
dividers = [5, 8, 11, 14, 17, 102]

rows_y = [round(1.42 + i * 0.57, 3) for i in range(7)]
for i, (lid, did) in enumerate(zip(labels, descs)):
    lab, d = now[lid], now[did]
    lab.left, lab.top, lab.width, lab.height = Inches(0.80), Inches(rows_y[i]), Inches(3.00), Inches(0.52)
    d.left,   d.top,   d.width,   d.height   = Inches(3.90), Inches(rows_y[i]), Inches(5.50), Inches(0.52)
for j, dvid in enumerate(dividers):
    dv = now[dvid]
    dv.left, dv.top, dv.width, dv.height = Inches(0.80), Inches(rows_y[j + 1] - 0.07), Inches(8.40), Inches(0.01)

# ----------------------------------------------------------------- SLIDE 5 (new)
layout = prs.slides[3].slide_layout
s5 = prs.slides.add_slide(layout)
for ph in list(s5.placeholders):
    ph._element.getparent().remove(ph._element)

def txt(x, y, w, h, s, size, bold, color, font="Calibri", align=PP_ALIGN.LEFT, anchor=None):
    tb = s5.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = s
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tb

txt(0.80, 0.50, 8.40, 0.70, TITLE5, 32, True, NIGHT, font="Cambria")
txt(0.80, 1.30, 8.40, 0.62, SUB5, 14, False, GRAY)

LY = [2.20, 2.80, 3.40, 4.00]
for i, (lab, note) in enumerate(LEVERS):
    txt(0.80, LY[i], 4.90, 0.55, lab, 15, False, NIGHT, anchor=MSO_ANCHOR.MIDDLE)
    txt(5.75, LY[i], 3.45, 0.55, note, 12, False, STONE, anchor=MSO_ANCHOR.MIDDLE)
for i in (1, 2, 3):
    y = LY[i] - 0.09
    conn = s5.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.80), Inches(y), Inches(9.00), Inches(y))
    conn.line.color.rgb = LINE; conn.line.width = Pt(0.75)

box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.70), Inches(4.75), Inches(8.60), Inches(0.62))
box.fill.solid(); box.fill.fore_color.rgb = SAND
box.line.fill.background()
box.shadow.inherit = False
tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = BANNER5
r.font.name = "Calibri"; r.font.size = Pt(13.5); r.font.bold = True; r.font.color.rgb = NIGHT

prs.save(SRC)
print("Saved. Slides now:", len(prs.slides.__iter__.__self__._sldIdLst))

# ------ verify ------
v = Presentation(SRC)
for si in (1, 4):
    s = v.slides[si]
    print(f"\n--- verify slide {si+1} ---")
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            from pptx.util import Emu
            print(f"  y={Emu(sh.top).inches:.2f} : {sh.text_frame.text[:70]!r}")
